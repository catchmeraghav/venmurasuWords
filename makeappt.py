#!/usr/bin/python
# -*- coding: utf-8 -*-


import os, re
import datetime
import random
from tamil import utf8
from downloadVenmurasu import downloadVenmurasu
from parseHTMLToTamilTxt import parseHTMLToTamilTxt
from makeSentences import makeSentencesForChapter
#from durusCommit import durusCommit

import pptx
import scipy.misc

class pptMaker(object):
    def __init__(self):
        pass
        self.prs = pptx.Presentation()
        
    def addAPicSlide(self, img_path, titleTxt=None, subtitleTxtDict=None):
        self.prs.slide_height = 5143500
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        pic_left  = int(0)
        pic_top   = int(0)
        if img_path:
            img = scipy.misc.imread(img_path)
            pic_width = None if int(self.prs.slide_width) > int(img.shape[1]) else int(self.prs.slide_width) 
            pic_height = int(img.shape[0]) if int(self.prs.slide_height) > int(img.shape[0]) else int(self.prs.slide_height)
            slide.shapes.add_picture(img_path, pic_left, pic_top, int(self.prs.slide_width), int(self.prs.slide_height))
            
            body_shape = slide.shapes.add_textbox(0, 0, int(self.prs.slide_width), int(self.prs.slide_height))
            removeSplChr = re.compile("[\"\'\[\]…:.-_*,;]‘’")
            
            if titleTxt:
                try:
                    titleTxt = removeSplChr.sub('', titleTxt)
                    titleTxt = ''.join(utf8.get_letters(titleTxt))
                    tf = body_shape.text_frame
                    tf.text = titleTxt
                except:
                    print "-- Exception at tile of  "
                    
            if subtitleTxtDict:
                def selectsub():
                    sentences = subtitleTxtDict
                    num = sentences.keys()[random.randrange(len(sentences.keys()))]
                    sentencesNumList = sentences.get(num)
                    return sentencesNumList[random.randrange(len(sentencesNumList))]
                shouldBreak = False
                subtitleTxt = None
                while not shouldBreak:
                    try:
                        subtitleTxt = selectsub()
                        subtitleTxt = unicode(subtitleTxt, "utf-8")     
                        subtitleTxt = removeSplChr.sub('', subtitleTxt)
                        subtitleTxt = ''.join(utf8.get_letters(subtitleTxt))
                        subtitleTxt = subtitleTxt.split(' ')
                        subtitleTxt = '\n'.join([' '.join(subtitleTxt[5 * i: 5 * i + 5]) for i in range(0, len(subtitleTxt) / 5)])
    
                        tf = body_shape.text_frame
                        p = tf.add_paragraph()
                        p.text = subtitleTxt
                        shouldBreak = True
                    except Exception, e:
                        #print "-- Exception at subtitle ", subtitleTxt,"\n",e
                        shouldBreak = False

    def saveppt(self, flname):
        self.prs.save('C:\\Users\\Rajasimhan\\Desktop\\venbooks\\complete-ven'+'.pptx')


def generateForDate():
    pass

dl = downloadVenmurasu()
def getStartDate(book):
    aDate = dl.books[book][0]
    return datetime.date(aDate[0], aDate[1], aDate[2])

def getEndDate(book):
    aDate = dl.books[book][1]
    return datetime.date(aDate[0], aDate[1], aDate[2])

def makeChapterDict(currentdate, folder):
    mkWds = makeSentencesForChapter(currentdate, folder)
    return mkWds.makeChapterDict()

def commitChapterWord(connect, currentdate, aKey, chapterDict):
    connect.makeChapterCommit(currentdate, aKey, chapterDict)


def completeGeneratePPT():
    out = pptMaker()
    try:
        
        for book in dl.books.keys():
            
            folder = os.path.join(dl.venmurasuFolder, book)
            
            if not os.path.exists(folder):
                os.makedirs(folder, mode=0777)
            folder = folder+"/"
            connect = durusCommit(folder, book)
            currentdate = getStartDate(book)
            enddate = getEndDate(book)
            print 'Folder name: ',folder
            while currentdate <= enddate:
                #print currentdate
                chapterSentenceDict = ''
                checkRoot = connect.getRoot()
                dateTuple = tuple(currentdate.strftime('%Y-%m-%d').split('-'))
                if checkRoot.get(book) and dateTuple in checkRoot[book]:
                    chapterSentenceDict = checkRoot[book][dateTuple]
                else:
                    chapterSentenceDict = makeChapterDict(currentdate, folder)
                if not isinstance(chapterSentenceDict, str):
                    commitChapterWord(connect, currentdate, book, chapterSentenceDict)
                    
                    imgName = currentdate.strftime('%Y-%m-%d')+'.jpg'
                    normalImgName = os.path.join(folder,imgName)
                    bigImgName = os.path.join(folder,'big-'+imgName)
                    
                    img_path = None
                    if bigImgName and os.path.isfile(bigImgName):
                        os.rename(bigImgName, normalImgName)
                        img_path = normalImgName
                    elif normalImgName and os.path.isfile(normalImgName):
                        img_path = normalImgName
                    
                    
                    subtitleTxt = chapterSentenceDict.get('sentenceDict')
                    out.addAPicSlide(img_path, titleTxt=chapterSentenceDict.get('title'), subtitleTxtDict=subtitleTxt)
                
                currentdate += datetime.timedelta(days=1)
            #out.saveppt(os.path.join(folder, book))
            print book," -- completed"
    except Exception, e:
        print e
    finally:
        out.saveppt(folder+'/complete')

def theLeftOutImg():
    dl = downloadVenmurasu()
    for book in dl.books.keys():
        folder = dl.venmurasuFolder+"/"+book
        
        if not os.path.exists(folder):
            os.makedirs(folder, mode=0777)
        folder = folder+"/"
        
        currentdate = getStartDate(book)
        enddate = getEndDate(book)
        print 'Folder name: ',folder
        while currentdate <= enddate:
            print currentdate
            
            imgName = currentdate.strftime('%Y-%m-%d')+'.jpg'
            normalImgName = os.path.join(folder,imgName)
            bigImgName = os.path.join(folder,'big-'+imgName)
                    
            if not (os.path.exists(bigImgName) or os.path.exists(normalImgName)):
                dl.downloadChapterImage(currentdate, folder)
                #print currentdate, "Image downloaded"
            currentdate += datetime.timedelta(days=1)
            
if __name__ == "__main__":
    #theLeftOutImg()
    completeGeneratePPT()

